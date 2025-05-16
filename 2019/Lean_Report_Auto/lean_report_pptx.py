import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.dml.color import RGBColor
#from lean_report_excel import LeanReportExcel
import configparser
import pandas as pd

#提案改善目标
def slider_goal(prs,propose_goal,save_manpower_goal,propose_actual,save_manpower_actual):
    #新建一页幻灯片
    slide=prs.slides.add_slide(prs.slide_layouts[1])
    #设定标题
    title_shape=slide.shapes.title
    title_shape.text="提案改善目标"
    #正文设定
    body_shape= slide.shapes.placeholders
    body_shape[1].text= '提案件数'
    
    new_paragraph= body_shape[1].text_frame.add_paragraph()
    new_paragraph.text= ' '#新段落中的文字
    new_paragraph.level = 0 #新段落的级别

    new_paragraph= body_shape[1].text_frame.add_paragraph()
    new_paragraph.text= ' '#新段落中的文字
    new_paragraph.level = 0 #新段落的级别

    new_paragraph= body_shape[1].text_frame.add_paragraph()
    new_paragraph.text= '节省人力'#新段落中的文字
    new_paragraph.level = 0 #新段落的级别

    #插入表格（提案状况）
    month_cn=["1月","2月","3月","4月","5月","6月","7月","8月","9月","10月","11月","12月"]
    
    rows, cols, left, top, width, height= 4, 14, Inches(0.5), Inches(2), Inches(9), Inches(0.8)
    table= slide.shapes.add_table(rows, cols, left, top, width, height).table #添加表格，并取表格类
    
    table.columns[0].width=Inches(1.2) #第一纵列宽度
    table.cell(1,0).text="目标提案件数"
    table.cell(1,0).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(2,0).text="实际提案件数"
    table.cell(2,0).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(3,0).text="提案达成率"
    table.cell(3,0).text_frame.paragraphs[0].font.size=Pt(10)

    for index in range(len(month_cn)):
        table.cell(0,index+1).text=month_cn[index]
        table.cell(0,index+1).text_frame.paragraphs[0].font.size=Pt(10)

    for index in range(len(propose_actual)):
        table.columns[index+1].width=Inches(0.6) #每个纵列宽度
        table.cell(1,index+1).text=str(propose_goal)
        table.cell(1,index+1).text_frame.paragraphs[0].font.size=Pt(10)
        table.cell(2,index+1).text=str(propose_actual[index+1])
        table.cell(2,index+1).text_frame.paragraphs[0].font.size=Pt(10)
        table.cell(3,index+1).text=str('{:.0%}'.format(int(propose_actual[index+1])/int(propose_goal)))
        table.cell(3,index+1).text_frame.paragraphs[0].font.size=Pt(10)
        if int(propose_actual[index+1])/int(propose_goal)>=1:
            table.cell(3,index+1).fill.solid()
            table.cell(3,index+1).fill.fore_color.rgb=RGBColor(0,255,0) 
        else:
            table.cell(3,index+1).fill.solid()
            table.cell(3,index+1).fill.fore_color.rgb=RGBColor(255,255,0) 

    table.columns[13].width=Inches(0.6) #最后一个纵列宽度
    table.cell(0,13).text="Total"
    table.cell(0,13).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(1,13).text=str(len(propose_actual)*int(propose_goal))
    table.cell(1,13).text_frame.paragraphs[0].font.size=Pt(10)
    total=0
    for index in range(len(propose_actual)):
        total=total+int(propose_actual[index+1]) 
    table.cell(2,13).text=str(total)
    table.cell(2,13).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(3,13).text=str('{:.0%}'.format(total/(len(propose_actual)*int(propose_goal))))
    table.cell(3,13).text_frame.paragraphs[0].font.size=Pt(10)
    if total/(len(propose_actual)*int(propose_goal))>=1:
        table.cell(3,13).fill.solid()
        table.cell(3,13).fill.fore_color.rgb=RGBColor(0,255,0) 
    else:
        table.cell(3,13).fill.solid()
        table.cell(3,13).fill.fore_color.rgb=RGBColor(255,255,0) 
    
    #插入表格（节省人力状况）
    month_cn=["1月","2月","3月","4月","5月","6月","7月","8月","9月","10月","11月","12月"]
    
    rows, cols, left, top, width, height= 4, 14, Inches(0.5), Inches(3.8), Inches(9), Inches(0.8)
    table= slide.shapes.add_table(rows, cols, left, top, width, height).table #添加表格，并取表格类
    
    table.columns[0].width=Inches(1.2) #第一纵列宽度
    table.cell(1,0).text="目标节省人数"
    table.cell(1,0).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(2,0).text="实际节省人数"
    table.cell(2,0).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(3,0).text="节省人力达成率"
    table.cell(3,0).text_frame.paragraphs[0].font.size=Pt(10)

    for index in range(len(month_cn)):
        table.cell(0,index+1).text=month_cn[index]
        table.cell(0,index+1).text_frame.paragraphs[0].font.size=Pt(10)

    for index in range(len(propose_actual)):
        table.columns[index+1].width=Inches(0.6) #每个纵列宽度
        table.cell(1,index+1).text=str(save_manpower_goal)
        table.cell(1,index+1).text_frame.paragraphs[0].font.size=Pt(10)
        table.cell(2,index+1).text=str(save_manpower_actual[index+1])
        table.cell(2,index+1).text_frame.paragraphs[0].font.size=Pt(10)
        table.cell(3,index+1).text=str('{:.0%}'.format(int(save_manpower_actual[index+1])/int(save_manpower_goal)))
        table.cell(3,index+1).text_frame.paragraphs[0].font.size=Pt(10)
        if int(save_manpower_actual[index+1])/int(save_manpower_goal)>=1:
            table.cell(3,index+1).fill.solid()
            table.cell(3,index+1).fill.fore_color.rgb=RGBColor(0,255,0) 
        else:
            table.cell(3,index+1).fill.solid()
            table.cell(3,index+1).fill.fore_color.rgb=RGBColor(255,255,0) 

    table.columns[13].width=Inches(0.6) #最后一个纵列宽度
    table.cell(0,13).text="Total"
    table.cell(0,13).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(1,13).text=str(len(propose_actual)*int(save_manpower_goal))
    table.cell(1,13).text_frame.paragraphs[0].font.size=Pt(10)
    total=0
    for index in range(len(propose_actual)):
        total=total+int(save_manpower_actual[index+1]) 
    table.cell(2,13).text=str(total)
    table.cell(2,13).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(3,13).text=str('{:.0%}'.format(total/(len(propose_actual)*int(save_manpower_goal))))
    table.cell(3,13).text_frame.paragraphs[0].font.size=Pt(10)
    if total/(len(propose_actual)*int(save_manpower_goal))>=1:
        table.cell(3,13).fill.solid()
        table.cell(3,13).fill.fore_color.rgb=RGBColor(0,255,0) 
    else:
        table.cell(3,13).fill.solid()
        table.cell(3,13).fill.fore_color.rgb=RGBColor(255,255,0)    

    return(prs)

#提案改善进度
def slider_implement(prs,data,subtitle):
    #新建一页幻灯片
    slide=prs.slides.add_slide(prs.slide_layouts[1])
    #设定标题
    title_shape=slide.shapes.title
    title_shape.text="提案改善进度"
    #正文设定
    body_shape= slide.shapes.placeholders
    body_shape[1].text= subtitle+"*"+str(len(data))
    #print(data.values[0][0])
    #print(data.columns.values)
    #制作表格
    rows, cols, left, top, width, height= len(data)+1, 7, Inches(0.1), Inches(2), Inches(9), Inches(0.8)
    table= slide.shapes.add_table(rows, cols, left, top, width, height).table #添加表格，并取表格类
    
    table.columns[0].width=Inches(0.3)
    table.columns[1].width=Inches(3.3)
    table.columns[2].width=Inches(1.2)
    table.columns[3].width=Inches(1.2)
    table.columns[4].width=Inches(0.8)
    table.columns[5].width=Inches(1)
    table.columns[6].width=Inches(2)
    
    table.cell(0,6).text="进度"
    table.cell(0,6).text_frame.paragraphs[0].font.size=Pt(10)

    for index in range(len(data.columns.values)):
        if index == 0:
            table.cell(0,index).text="项次"
            table.cell(0,index).text_frame.paragraphs[0].font.size=Pt(10)
        else:
            table.cell(0,index).text=data.columns.values[index]
            table.cell(0,index).text_frame.paragraphs[0].font.size=Pt(10)

    for row in range(len(data.values)):
        for col in range(6):
            if ".xml" in str(data.values[row][col]):
                table.cell(row+1,col).text=str(row+1)
                table.cell(row+1,col).text_frame.paragraphs[0].font.size=Pt(10)
            elif "/WKS/Wistron" in str(data.values[row][col]):
                table.cell(row+1,col).text=data.values[row][col].split("/")[0]
                table.cell(row+1,col).text_frame.paragraphs[0].font.size=Pt(10)
            else:
                table.cell(row+1,col).text=str(data.values[row][col])
                table.cell(row+1,col).text_frame.paragraphs[0].font.size=Pt(10)
           
    #print(pd.to_datetime(data['預計完成日期'],format="%m/%d/%y"))
    #print(data['預計完成日期'].apply(lambda x: x.strftime('%Y-%m-%d')))
    #print(data['預計完成日期'].dt.strftime('%m/%d/%Y').values)
    for index in range(len(data['預計完成日期'].dt.strftime('%m/%d/%Y').values)):
        table.cell(index+1,5).text=data['預計完成日期'].dt.strftime('%m/%d/%Y').values[index]
        table.cell(index+1,5).text_frame.paragraphs[0].font.size=Pt(10)
    
    return(prs)

#专案改善进度
def slider_project(prs,data):
    #新建一页幻灯片
    slide=prs.slides.add_slide(prs.slide_layouts[1])
    #设定标题
    title_shape=slide.shapes.title
    title_shape.text="专案改善进度"
    #正文设定
    body_shape= slide.shapes.placeholders
    body_shape[1].text= "专案改善Total*"+str(len(data))

    #制作表格
    '''rows, cols, left, top, width, height= len(data)+1, 6, Inches(0.1), Inches(2), Inches(9), Inches(0.8)
    table= slide.shapes.add_table(rows, cols, left, top, width, height).table #添加表格，并取表格类
    
    table.columns[0].width=Inches(0.5)
    table.columns[1].width=Inches(5.8)
    table.columns[2].width=Inches(1)
    table.columns[3].width=Inches(0.8)
    table.columns[4].width=Inches(0.8)
    table.columns[5].width=Inches(0.9)

    table.cell(0,0).text="NO."
    table.cell(0,0).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(0,1).text="Items"
    table.cell(0,1).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(0,2).text="D.D."
    table.cell(0,2).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(0,3).text="PIC"
    table.cell(0,3).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(0,4).text="Status"
    table.cell(0,4).text_frame.paragraphs[0].font.size=Pt(10)
    table.cell(0,5).text="Saving"
    table.cell(0,5).text_frame.paragraphs[0].font.size=Pt(10)
    

    for row in range(len(data.values)):
        for col in range(len(data.values[0])):
            table.cell(row+1,col).text=str(data.values[row][col])
            for paras in table.cell(row+1,col).text_frame.paragraphs:
                paras.font.size=Pt(10)
            if str(data.values[row][col])=="Closed":
                table.cell(row+1,col).fill.solid()
                table.cell(row+1,col).fill.fore_color.rgb=RGBColor(0,255,0) 
            elif str(data.values[row][col])=="Ongoing":
                table.cell(row+1,col).fill.solid()
                table.cell(row+1,col).fill.fore_color.rgb=RGBColor(255,255,0)
            elif str(data.values[row][col])=="Pending":
                table.cell(row+1,col).fill.solid()
                table.cell(row+1,col).fill.fore_color.rgb=RGBColor(192,192,192)


    for index in range(len(data['D.D.'].dt.strftime('%m/%d/%Y').values)):
        table.cell(index+1,2).text=data['D.D.'].dt.strftime('%m/%d/%Y').values[index]
        table.cell(index+1,2).text_frame.paragraphs[0].font.size=Pt(10)'''

    return(prs)

def slider_reward(prs):
    #新建一页幻灯片
    slide=prs.slides.add_slide(prs.slide_layouts[1])
    #设定标题
    title_shape=slide.shapes.title
    title_shape.text="嘉奖&部费状况"
    #正文设定
    body_shape= slide.shapes.placeholders
    body_shape[1].text= "嘉奖状况"

    return(prs)

if __name__ == "__main__":
    propose_actual={1: 0}
    save_manpower_actual={1: 1, 2: 0, 3: 0, 4: 0, 5: 0, 6: 0, 7: 0, 8: 0, 9: 0, 10: 0, 11: 0, 12: 0}
    #获取ini档案信息
    config=configparser.ConfigParser()
    config.read(r"C:\Users\Charles\Desktop\Lean_Report_Auto\config.ini")
    propose_goal=config.get('info', 'propose_goal')
    save_manpower_goal=config.get('info', 'save_manpower_goal')
    #加载模板
    prs=Presentation(r'C:\Users\Charles\Desktop\Lean_Report_Auto\Lean Report Template.pptx')
    #制作第一页
    prs=slider_goal(prs,propose_goal,save_manpower_goal,propose_actual,save_manpower_actual)
    #制作第二页
    implement_te_data=pd.read_excel(r"C:\Users\Charles\Desktop\Lean_Report_Auto\implement_te.xlsx")
    prs=slider_implement(prs,implement_te_data,"待TE签核")
    #制作第三页
    implement_ie_data=pd.read_excel(r"C:\Users\Charles\Desktop\Lean_Report_Auto\implement_ie.xlsx")
    prs=slider_implement(prs,implement_ie_data,"待IE签核")
    #制作第四页
    project_data=pd.read_excel(r"C:\Users\Charles\Desktop\Lean_Report_Auto\專案改善進度.xlsx")
    prs=slider_project(prs,project_data)
    #制作第五页
    prs=slider_reward(prs)
    #获取当前日期
    today=datetime.date.today()
    #保存pptx
    prs.save(r'C:\Users\Charles\Desktop\Lean_Report_Auto\Lean Report %s.pptx' % (str(today)))


